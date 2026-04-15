"""Build Investor Deck, Board Deck, and Town Hall Deck from the cloned .potx master.

Pattern: same as skills/pptx — clone the .potx, rewrite content type to .pptx,
open with python-pptx, add slides using the NDLNG layouts, fill the placeholder
text (title + body), leave the brand chrome / logos / colors untouched.
"""
import sys, zipfile, shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

POTX = Path("/home/user/Skills-fork/uploads/05-presentation-templates/"
            "NextDecade Power Point Slide Master_Final_Oct 2025.potx")

# Layout index map (verified from the .potx)
LAYOUT = {
    "cover":          0,   # 'Custom Layout'
    "disclaimer":     2,   # 'Public Disclaimer'
    "back_cover":     3,   # '23_Custom Layout'
    "nd_blank":       11,
    "nd_single":      12,  # 'ND Single Narrative 18 Font'
    "nd_dual":        13,
    "nd_tri_right":   14,
    "nd_tri_left":    15,
    "nd_quad":        16,
    "rg_blank":       23,
    "rg_single":      24,
    "rg_dual":        25,
    "rg_tri_left":    26,
    "rg_tri_right":   27,
    "rg_quad":        28,
    "ncs_blank":      33,
    "ncs_single":     34,
}

NAVY = RGBColor(0x00, 0x20, 0x60)
ORANGE = RGBColor(0xFC, 0x71, 0x34)


def clone_potx_as_pptx(out_pptx: Path):
    """Clone the .potx to a .pptx by rewriting the content type."""
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(str(POTX), "r") as zin:
        with zipfile.ZipFile(str(out_pptx), "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(
                        b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                        b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                    )
                zout.writestr(item, data)


def _remove_all_slides(prs: Presentation):
    """Remove any pre-existing slides (.potx may or may not ship with any)."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _fill_title(slide, text: str):
    if slide.shapes.title is not None:
        slide.shapes.title.text = text


def _fill_by_idx(slide, idx: int, text: str):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return


def _add_content_bullets(slide, bullets: list[str], body_idx: int = 10):
    """Replace the content placeholder text with bullets."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == body_idx and ph.has_text_frame:
            tf = ph.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0
            return


def _add_slide(prs, layout_idx: int):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


# ---------- INVESTOR DECK ----------

def build_investor_deck(output: Path):
    clone_potx_as_pptx(output)
    prs = Presentation(str(output))
    _remove_all_slides(prs)

    # 1. Cover
    s = _add_slide(prs, LAYOUT["cover"])
    _fill_title(s, "Investor Update")
    _fill_by_idx(s, 10, "April 15, 2026")

    # 2. Public Disclaimer
    _add_slide(prs, LAYOUT["disclaimer"])  # layout owns the FLS legal text

    # 3. Agenda
    s = _add_slide(prs, LAYOUT["nd_single"])
    _fill_title(s, "Agenda")
    _add_content_bullets(s, [
        "Business highlights - Q1 2026",
        "Rio Grande LNG Phase 1 construction status",
        "Commercial and regulatory update",
        "Financial position and liquidity",
        "Phase 2 path to Final Investment Decision",
        "Q&A",
    ])

    # 4. RG Phase 1 progress
    s = _add_slide(prs, LAYOUT["rg_single"])
    _fill_title(s, "Rio Grande LNG Phase 1 - 68% Complete")
    _add_content_bullets(s, [
        "Engineering 99% complete",
        "Procurement 88% complete - all long-lead equipment on site or in transit",
        "Construction 52% complete across Trains 1, 2, and 3",
        "Train 1 cold box set during Q1; no incidents",
        "First LNG target unchanged: second half of 2027",
    ])

    # 5. Safety
    s = _add_slide(prs, LAYOUT["rg_single"])
    _fill_title(s, "Safety Performance")
    _add_content_bullets(s, [
        "3.1 million construction hours in Q1 2026",
        "Total Recordable Incident Rate: 0.42 (industry benchmark 0.65)",
        "No Life-Saving Rule violations recorded",
        "Near-miss event HW-2026-0143 investigated; corrective actions complete",
        "Life-Saving Rules stand-down completed site-wide on April 14, 2026",
    ])

    # 6. Commercial
    s = _add_slide(prs, LAYOUT["nd_single"])
    _fill_title(s, "Phase 2 Commercial Progress")
    _add_content_bullets(s, [
        "Three investment-grade counterparties in advanced SPA negotiations",
        "FERC Phase 2 authorization remains in effect",
        "DOE non-FTA export authorization expected Q2 2026",
        "Targeting a Final Investment Decision on Phase 2 in the second half of 2026",
    ])

    # 7. NCS
    s = _add_slide(prs, LAYOUT["ncs_single"])
    _fill_title(s, "NEXT Carbon Solutions - CCS Pathway")
    _add_content_bullets(s, [
        "CCS project designed to permanently store CO2 from liquefaction",
        "Site characterization underway; Class VI permits advancing",
        "Supports Rio Grande LNG's lower carbon intensive value proposition",
        "First injection targeted after Phase 1 first LNG",
    ])

    # 8. Financial
    s = _add_slide(prs, LAYOUT["nd_single"])
    _fill_title(s, "Financial Summary")
    _add_content_bullets(s, [
        "Phase 1 budget unchanged at $18.4B",
        "Q1 2026 total project spend: $9.8B",
        "Funding mix: project debt $6.2B, committed equity $3.0B, revolver $0.6B",
        "Liquidity adequate through first LNG",
    ])

    # 9. Closing
    s = _add_slide(prs, LAYOUT["nd_single"])
    _fill_title(s, "Closing and Q&A")
    _add_content_bullets(s, [
        "Construction on schedule; first LNG 2H 2027",
        "Phase 2 FID targeted 2H 2026",
        "Safety performance ahead of benchmark",
        "ir@next-decade.com",
    ])

    # 10. Back cover
    _add_slide(prs, LAYOUT["back_cover"])

    prs.save(str(output))
    print(f"Wrote {output}")


# ---------- BOARD DECK ----------

def build_board_deck(output: Path):
    clone_potx_as_pptx(output)
    prs = Presentation(str(output))
    _remove_all_slides(prs)

    # 1. Cover
    s = _add_slide(prs, LAYOUT["cover"])
    _fill_title(s, "Q1 2026 Board Update")
    _fill_by_idx(s, 10, "April 15, 2026")

    # 2. Public Disclaimer (board decks also ship with disclaimer when externally archived)
    _add_slide(prs, LAYOUT["disclaimer"])

    # 3. Agenda
    s = _add_slide(prs, LAYOUT["nd_single"])
    _fill_title(s, "Agenda")
    _add_content_bullets(s, [
        "Safety moment",
        "CEO highlights",
        "Construction progress (COO)",
        "Financial summary (CFO)",
        "Commercial update (VP Commercial)",
        "Phase 2 FID readiness",
        "Executive session",
    ])

    # 4. Safety moment
    s = _add_slide(prs, LAYOUT["rg_single"])
    _fill_title(s, "Safety Moment: Hot Work and Fire Watch")
    _add_content_bullets(s, [
        "HW-2026-0143 near-miss on April 14, 2026",
        "Root cause: fire watch relieved without formal handoff",
        "Corrective actions complete",
        "HSSE Committee to review at May 2026 session",
    ])

    # 5. CEO highlights
    s = _add_slide(prs, LAYOUT["nd_single"])
    _fill_title(s, "CEO Highlights - Q1 2026")
    _add_content_bullets(s, [
        "Phase 1 construction 68% complete",
        "First LNG target 2H 2027 unchanged",
        "Phase 2 SPA negotiations advancing with three investment-grade counterparties",
        "No material safety events recorded",
    ])

    # 6. Construction
    s = _add_slide(prs, LAYOUT["rg_single"])
    _fill_title(s, "Construction Progress")
    _add_content_bullets(s, [
        "Engineering 99%, procurement 88%, construction 52%",
        "Train 1 cold box set in Q1 without incident",
        "Non-cryogenic utilities commissioning begins July 2026",
        "Schedule confidence: high",
    ])

    # 7. Financial
    s = _add_slide(prs, LAYOUT["nd_single"])
    _fill_title(s, "Financial Summary")
    _add_content_bullets(s, [
        "Phase 1 remains on budget ($18.4B)",
        "Q1 spend: $9.8B; adequate liquidity through first LNG",
        "Phase 2 financing roadmap on track for 2H 2026 FID package",
    ])

    # 8. Phase 2 FID
    s = _add_slide(prs, LAYOUT["nd_single"])
    _fill_title(s, "Phase 2 FID Readiness")
    _add_content_bullets(s, [
        "Pre-FID items: SPA execution, debt commitments, EPC cost certainty, equity funding, community engagement review",
        "Board action requested: authorize management to prepare the Phase 2 FID package for August 2026 meeting",
    ])

    # 9. Executive session
    s = _add_slide(prs, LAYOUT["nd_blank"])
    _fill_title(s, "Executive Session")

    # 10. Back cover
    _add_slide(prs, LAYOUT["back_cover"])

    prs.save(str(output))
    print(f"Wrote {output}")


# ---------- TOWN HALL DECK (casual internal voice) ----------

def build_town_hall_deck(output: Path):
    clone_potx_as_pptx(output)
    prs = Presentation(str(output))
    _remove_all_slides(prs)

    # 1. Cover (no disclaimer - internal deck)
    s = _add_slide(prs, LAYOUT["cover"])
    _fill_title(s, "All-Hands Town Hall")
    _fill_by_idx(s, 10, "April 15, 2026")

    # Agenda
    s = _add_slide(prs, LAYOUT["nd_single"])
    _fill_title(s, "What we'll cover")
    _add_content_bullets(s, [
        "Where we are - Phase 1 construction",
        "Our safety performance",
        "Phase 2 - what's next",
        "People news",
        "Questions - we've got time for you",
    ])

    # Construction
    s = _add_slide(prs, LAYOUT["rg_single"])
    _fill_title(s, "Construction: we're 68% there")
    _add_content_bullets(s, [
        "That's a huge milestone - thank you to everyone who got us here",
        "Engineering is basically done; procurement is almost there",
        "Train 1 cold box set with zero incidents - incredible work",
        "First LNG is still on track for the second half of 2027",
    ])

    # Safety
    s = _add_slide(prs, LAYOUT["rg_single"])
    _fill_title(s, "Safety - we're doing well, and we can be better")
    _add_content_bullets(s, [
        "3.1 million hours this quarter - TRIR 0.42, better than the industry",
        "We had a hot-work near-miss on April 14 - HW-2026-0143",
        "The fire watch relief step wasn't formally handed off. We closed the gap.",
        "Reminder: every permit gets a formal handoff, every time",
    ])

    # Phase 2
    s = _add_slide(prs, LAYOUT["nd_single"])
    _fill_title(s, "Phase 2 - what's next")
    _add_content_bullets(s, [
        "Commercial team is advancing three long-term SPAs",
        "FERC authorization is in place; DOE is expected in Q2",
        "Our goal is a Phase 2 Final Investment Decision in the second half of 2026",
        "We can't do it without every one of you",
    ])

    # People
    s = _add_slide(prs, LAYOUT["nd_single"])
    _fill_title(s, "People news")
    _add_content_bullets(s, [
        "47 new hires in Q1 - welcome!",
        "Two teammates celebrated 5-year anniversaries this month",
        "NEXT Digest newsletter now has a 'day in the life' section - submit yours",
        "Benefits open enrollment opens in May",
    ])

    # Close
    s = _add_slide(prs, LAYOUT["nd_single"])
    _fill_title(s, "Over to you")
    _add_content_bullets(s, [
        "Questions - raise your hand or drop them in the chat",
        "Can't stay? Send us a note at allhands@next-decade.com",
        "Thanks, team",
    ])

    # Back cover
    _add_slide(prs, LAYOUT["back_cover"])

    prs.save(str(output))
    print(f"Wrote {output}")


if __name__ == "__main__":
    out_dir = Path(sys.argv[1])
    out_dir.mkdir(parents=True, exist_ok=True)
    build_investor_deck(out_dir / "Investor Update - Q1 2026.pptx")
    build_board_deck(out_dir / "Q1 2026 Board Deck.pptx")
    build_town_hall_deck(out_dir / "All-Hands Town Hall - April 2026.pptx")
